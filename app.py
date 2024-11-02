from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from transformers import pipeline
import os

app = Flask(__name__)

# Load the summarization model locally
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# Route for the homepage to serve the HTML form
@app.route('/')
def index():
    return render_template("index.html")

def generate_summary(text):
    max_input_length = 3000  # Define maximum input length
    max_chunk_length = 1024   # Define max chunk length for model (depends on model capacity)

    # Split the text into chunks if it's longer than max_input_length
    if len(text.split()) > max_input_length:
        words = text.split()
        chunks = [' '.join(words[i:i + max_chunk_length]) for i in range(0, len(words), max_chunk_length)]
    else:
        chunks = [text]

    # Generate summaries for each chunk
    summary_chunks = []
    for chunk in chunks:
        summary = summarizer(chunk, max_length=130, min_length=30, do_sample=False)
        summary_chunks.append(summary[0]['summary_text'])

    # Combine summaries
    full_summary = ' '.join(summary_chunks)
    return full_summary

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_ppt(slides_content, title):
    prs = Presentation()

    # Define the RGB color for grey background
    grey_color = RGBColor(169, 169, 169)  # Light grey color

    for slide_text in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using title and content layout

        # Set grey background
        slide_background = slide.background
        fill = slide_background.fill
        fill.solid()
        fill.fore_color.rgb = grey_color

        title_placeholder = slide.shapes.title
        content_placeholder = slide.shapes.placeholders[1]

        # Set title text
        title_placeholder.text = title
        title_placeholder.text_frame.paragraphs[0].alignment = 1  # Center alignment for title

        # Set content text
        content_placeholder.text = slide_text
        for paragraph in content_placeholder.text_frame.paragraphs:
            paragraph.alignment = 1  # Center alignment for content
            for run in paragraph.runs:
                run.font.size = Pt(18)  # Optional: Set font size for better readability

    # Save the presentation in the current directory
    pptx_path = os.path.join(os.getcwd(), 'generated_presentation.pptx')
    
    prs.save(pptx_path)
    
    return pptx_path

@app.route('/convert', methods=['POST'])
def convert_text_to_ppt():
    data = request.json
    text = data.get('text')
    title = data.get('title', 'Generated Presentation')
    
    if not text:
        return jsonify({"error": "No text provided"}), 400

    summary = generate_summary(text)
    slides_content = summary.split(". ")

    ppt_file_path = create_ppt(slides_content, title)
    return send_file(ppt_file_path, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)